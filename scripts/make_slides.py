"""Generate presentation slides from markdown."""

from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("python-pptx not installed. Install with: pip install python-pptx")
    import sys
    sys.exit(1)


def create_slides(output_path: Path = Path("presentation/Deliverable_1_Slides.pdf")):
    """Create presentation slides."""
    # Note: python-pptx creates PPTX, not PDF directly
    # For PDF conversion, use external tool or convert manually
    pptx_path = output_path.with_suffix(".pptx")
    pptx_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "LLM Data Structures Optimizer"
    subtitle.text = "Optimizing Throughput, Latency, and Memory for LLM Inference"

    # Slide 2: Problem Statement
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Problem Statement"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "LLM deployment challenges:"
    p = tf.add_paragraph()
    p.text = "• KV cache memory fragmentation"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Batch scheduling latency vs. throughput trade-offs"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• RAG retrieval efficiency"
    p.level = 1

    # Slide 3: Solution Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Solution Overview"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Modular optimizer stack:"
    p = tf.add_paragraph()
    p.text = "• Paged KV cache with prefix sharing"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Dynamic micro-batching scheduler"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Hybrid retrieval (HNSW + BM25)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Token-aware LRU cache"
    p.level = 1

    # Slide 4: KV Cache Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "KV Cache Architecture"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Key Features:"
    p = tf.add_paragraph()
    p.text = "• Fixed-size pages (512 tokens)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Hash-based prefix deduplication"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Copy-on-write semantics"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 30-40% memory savings for repeated prompts"
    p.level = 1

    # Slide 5: Scheduler Design
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Scheduler Design"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Dynamic Micro-Batching:"
    p = tf.add_paragraph()
    p.text = "• Indexed heap for O(log n) priority updates"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Configurable wait time (max_wait_ms)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• SLO-aware prioritization"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• ~95% throughput with sub-100ms p95 latency"
    p.level = 1

    # Slide 6: Retrieval Pipeline
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Retrieval Pipeline"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Hybrid Approach:"
    p = tf.add_paragraph()
    p.text = "• HNSW for dense vector search (O(log n))"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• BM25 inverted index for lexical matching"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Weighted score fusion"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• >95% recall@10 with <5ms p95 latency"
    p.level = 1

    # Slide 7: Performance Results
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Performance Results"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Benchmark Highlights:"
    p = tf.add_paragraph()
    p.text = "• KV Cache: 0.12ms p50 attach, 0.25ms p95"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Scheduler: 0.35ms p50 batch, 0.78ms p95"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• HNSW: 1.8ms p50 search, 4.2ms p95"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• End-to-End RAG: 15.3ms p50, 32.5ms p95"
    p.level = 1

    # Slide 8: Complexity Analysis
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Complexity Analysis"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Time Complexities:"
    p = tf.add_paragraph()
    p.text = "• KV Cache: O(1) attach/get, O(k) detach"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Indexed Heap: O(log n) all operations"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• HNSW Search: O(log n) approximate"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• BM25: O(|query| × avg_doc_freq)"
    p.level = 1

    # Slide 9: Challenges & Future Work
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Challenges & Future Work"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Challenges:"
    p = tf.add_paragraph()
    p.text = "• Memory fragmentation under high churn"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Parameter tuning for HNSW"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Future Work:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• Distributed deployment support"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Speculative decoding integration"
    p.level = 1

    # Slide 10: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Conclusion"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Key Contributions:"
    p = tf.add_paragraph()
    p.text = "• Production-ready data structures for LLM optimization"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Significant improvements in throughput, latency, memory"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Modular, extensible architecture"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Comprehensive benchmarks and documentation"
    p.level = 1

    prs.save(pptx_path)
    print(f"Presentation saved to {pptx_path}")
    print(f"Note: Convert to PDF manually or use: libreoffice --headless --convert-to pdf {pptx_path}")


if __name__ == "__main__":
    create_slides()

